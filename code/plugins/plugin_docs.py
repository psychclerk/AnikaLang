import os
import subprocess

from .base_plugin import AnikaPlugin
from core.interpreter import NativeFunction
from core.errors import FMS_Error
from core.utils import _register_doc, _get_doc, _close_doc

class DocsPlugin(AnikaPlugin):
    @staticmethod
    def _get_docx():
        try:
            import docx
            return docx
        except ImportError:
            raise FMS_Error("DOCX support requires python-docx. Run: pip install python-docx", error_type="Import Error")

    @staticmethod
    def _get_pptx():
        try:
            import pptx
            return pptx
        except ImportError:
            raise FMS_Error("PPTX support requires python-pptx. Run: pip install python-pptx", error_type="Import Error")

    def register(self, env, interpreter):
        # ==========================================================================
        # DOCX (Word Document) MANIPULATION
        # ==========================================================================
        def docx_create(i, a):
            docx = self._get_docx()
            doc = docx.Document()
            return _register_doc(doc)
        def docx_open(i, a):
            docx = self._get_docx()
            path = str(a[0])
            try:
                doc = docx.Document(path)
                return _register_doc(doc)
            except FileNotFoundError:
                raise FMS_Error(f"DOCX file not found: '{path}'", error_type="File Error")
            except Exception as e:
                raise FMS_Error(f"Failed to open DOCX: {str(e)}", error_type="File Error")
        def docx_save(i, a):
            handle, path = a[0], str(a[1])
            doc = _get_doc(handle)
            try:
                os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
                doc.save(path)
                return "SUCCESS"
            except Exception as e:
                raise FMS_Error(f"Failed to save DOCX: {str(e)}", error_type="File Error")
        def docx_close(i, a):
            _close_doc(a[0])
            return None
        def docx_add_paragraph(i, a):
            doc = _get_doc(a[0])
            text = str(a[1])
            style = str(a[2]) if len(a) > 2 and a[2] else None
            try:
                if style: doc.add_paragraph(text, style=style)
                else: doc.add_paragraph(text)
                return "SUCCESS"
            except Exception as e:
                raise FMS_Error(f"Failed to add paragraph: {str(e)}", error_type="Runtime Error")
        def docx_add_heading(i, a):
            doc = _get_doc(a[0])
            text = str(a[1])
            level = int(a[2]) if len(a) > 2 and a[2] else 1
            try:
                doc.add_heading(text, level=level)
                return "SUCCESS"
            except Exception as e:
                raise FMS_Error(f"Failed to add heading: {str(e)}", error_type="Runtime Error")
        def docx_add_bullet(i, a):
            doc = _get_doc(a[0])
            text = str(a[1])
            level = int(a[2]) if len(a) > 2 and a[2] else 0
            try:
                style_name = f'List Bullet {"1234"[level]}' if level > 0 else 'List Bullet'
                doc.add_paragraph(text, style=style_name)
                return "SUCCESS"
            except Exception:
                bullets = ["- "]
                prefix = bullets[min(level, len(bullets) - 1)]
                doc.add_paragraph(prefix + text)
                return "SUCCESS"
        def docx_add_numbered(i, a):
            doc = _get_doc(a[0])
            text = str(a[1])
            level = int(a[2]) if len(a) > 2 and a[2] else 0
            try:
                style_name = f'List Number {"1234"[level]}' if level > 0 else 'List Number'
                doc.add_paragraph(text, style=style_name)
                return "SUCCESS"
            except Exception:
                doc.add_paragraph(text)
                return "SUCCESS"
        def docx_add_page_break(i, a):
            doc = _get_doc(a[0])
            doc.add_page_break()
            return "SUCCESS"
        def docx_add_table(i, a):
            doc = _get_doc(a[0])
            data = a[1]
            has_header = bool(a[2]) if len(a) > 2 else True
            if not isinstance(data, list) or len(data) == 0:
                raise FMS_Error("DOCX_ADD_TABLE requires a non-empty list of lists", error_type="Runtime Error")
            try:
                rows = len(data)
                cols = max(len(row) for row in data)
                table = doc.add_table(rows=rows, cols=cols)
                table.style = 'Table Grid'
                for r, row_data in enumerate(data):
                    for c, cell_val in enumerate(row_data):
                        if c < cols:
                            table.cell(r, c).text = str(cell_val) if cell_val is not None else ""
                if has_header and len(data) > 0:
                    for cell in table.rows[0].cells:
                        for paragraph in cell.paragraphs:
                            for run in paragraph.runs:
                                run.font.bold = True
                return "SUCCESS"
            except Exception as e:
                raise FMS_Error(f"Failed to add table: {str(e)}", error_type="Runtime Error")
        def docx_add_image(i, a):
            doc = _get_doc(a[0])
            img_path = str(a[1])
            width = float(a[2]) if len(a) > 2 and a[2] else 4.0
            try:
                from docx.shared import Inches
                doc.add_picture(img_path, width=Inches(width))
                return "SUCCESS"
            except FileNotFoundError:
                raise FMS_Error(f"Image not found: '{img_path}'", error_type="File Error")
            except Exception as e:
                raise FMS_Error(f"Failed to add image: {str(e)}", error_type="Runtime Error")
        def docx_get_text(i, a):
            doc = _get_doc(a[0])
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return "\n".join(paragraphs)
        def docx_get_paragraphs(i, a):
            doc = _get_doc(a[0])
            return [para.text for para in doc.paragraphs]
        def docx_replace_text(i, a):
            doc = _get_doc(a[0])
            old_text, new_text = str(a[1]), str(a[2])
            count = 0
            for para in doc.paragraphs:
                if old_text in para.text:
                    for run in para.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            count += 1
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            if old_text in para.text:
                                for run in para.runs:
                                    if old_text in run.text:
                                        run.text = run.text.replace(old_text, new_text)
                                        count += 1
            return count
        def docx_set_header(i, a):
            doc = _get_doc(a[0])
            text = str(a[1])
            section = doc.sections[0]
            section.header.paragraphs[0].text = text
            return "SUCCESS"
        def docx_set_footer(i, a):
            doc = _get_doc(a[0])
            text = str(a[1])
            section = doc.sections[0]
            section.footer.paragraphs[0].text = text
            return "SUCCESS"
        def docx_set_title(i, a):
            doc = _get_doc(a[0])
            doc.core_properties.title = str(a[1])
            return "SUCCESS"
        def docx_set_author(i, a):
            doc = _get_doc(a[0])
            doc.core_properties.author = str(a[1])
            return "SUCCESS"
        def docx_get_metadata(i, a):
            doc = _get_doc(a[0])
            props = doc.core_properties
            return {
                "title": props.title or "", "author": props.author or "",
                "subject": props.subject or "", "created": str(props.created) if props.created else "",
                "modified": str(props.modified) if props.modified else "", "keywords": props.keywords or ""
            }
        def docx_to_text(i, a):
            docx = self._get_docx()
            docx_path, txt_path = str(a[0]), str(a[1])
            try:
                doc = docx.Document(docx_path)
                text = "\n".join([p.text for p in doc.paragraphs])
                with open(txt_path, 'w', encoding='utf-8') as f: f.write(text)
                return "SUCCESS"
            except Exception as e:
                raise FMS_Error(f"DOCX to text conversion failed: {str(e)}", error_type="Runtime Error")

        env.define("DOCX_CREATE", NativeFunction("DOCX_CREATE", 0, docx_create))
        env.define("DOCX_OPEN", NativeFunction("DOCX_OPEN", 1, docx_open))
        env.define("DOCX_SAVE", NativeFunction("DOCX_SAVE", 2, docx_save))
        env.define("DOCX_CLOSE", NativeFunction("DOCX_CLOSE", 1, docx_close))
        env.define("DOCX_ADD_PARAGRAPH", NativeFunction("DOCX_ADD_PARAGRAPH", -1, docx_add_paragraph))
        env.define("DOCX_ADD_HEADING", NativeFunction("DOCX_ADD_HEADING", -1, docx_add_heading))
        env.define("DOCX_ADD_BULLET", NativeFunction("DOCX_ADD_BULLET", -1, docx_add_bullet))
        env.define("DOCX_ADD_NUMBERED", NativeFunction("DOCX_ADD_NUMBERED", -1, docx_add_numbered))
        env.define("DOCX_ADD_PAGE_BREAK", NativeFunction("DOCX_ADD_PAGE_BREAK", 1, docx_add_page_break))
        env.define("DOCX_ADD_TABLE", NativeFunction("DOCX_ADD_TABLE", -1, docx_add_table))
        env.define("DOCX_ADD_IMAGE", NativeFunction("DOCX_ADD_IMAGE", -1, docx_add_image))
        env.define("DOCX_GET_TEXT", NativeFunction("DOCX_GET_TEXT", 1, docx_get_text))
        env.define("DOCX_GET_PARAGRAPHS", NativeFunction("DOCX_GET_PARAGRAPHS", 1, docx_get_paragraphs))
        env.define("DOCX_REPLACE_TEXT", NativeFunction("DOCX_REPLACE_TEXT", 3, docx_replace_text))
        env.define("DOCX_SET_HEADER", NativeFunction("DOCX_SET_HEADER", 2, docx_set_header))
        env.define("DOCX_SET_FOOTER", NativeFunction("DOCX_SET_FOOTER", 2, docx_set_footer))
        env.define("DOCX_SET_TITLE", NativeFunction("DOCX_SET_TITLE", 2, docx_set_title))
        env.define("DOCX_SET_AUTHOR", NativeFunction("DOCX_SET_AUTHOR", 2, docx_set_author))
        env.define("DOCX_GET_METADATA", NativeFunction("DOCX_GET_METADATA", 1, docx_get_metadata))
        env.define("DOCX_TO_TEXT", NativeFunction("DOCX_TO_TEXT", 2, docx_to_text))

        # ==========================================================================
        # PPTX (PowerPoint) MANIPULATION
        # ==========================================================================
        def pptx_create(i, a):
            pptx = self._get_pptx()
            prs = pptx.Presentation()
            return _register_doc(prs)
        def pptx_open(i, a):
            pptx = self._get_pptx()
            path = str(a[0])
            try:
                prs = pptx.Presentation(path)
                return _register_doc(prs)
            except FileNotFoundError:
                raise FMS_Error(f"PPTX file not found: '{path}'", error_type="File Error")
            except Exception as e:
                raise FMS_Error(f"Failed to open PPTX: {str(e)}", error_type="File Error")
        def pptx_save(i, a):
            handle, path = a[0], str(a[1])
            prs = _get_doc(handle)
            try:
                os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
                prs.save(path)
                return "SUCCESS"
            except Exception as e:
                raise FMS_Error(f"Failed to save PPTX: {str(e)}", error_type="File Error")
        def pptx_close(i, a):
            _close_doc(a[0])
            return None
        def pptx_add_title_slide(i, a):
            prs = _get_doc(a[0])
            title = str(a[1])
            subtitle = str(a[2]) if len(a) > 2 and a[2] else ""
            try:
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                if subtitle and len(slide.placeholders) > 1:
                    slide.placeholders[1].text = subtitle
                return len(prs.slides) - 1
            except Exception as e:
                raise FMS_Error(f"Failed to add title slide: {str(e)}", error_type="Runtime Error")
        def pptx_add_slide(i, a):
            prs = _get_doc(a[0])
            layout_idx = int(a[1]) if len(a) > 1 and a[1] is not None else 5
            try:
                if layout_idx >= len(prs.slide_layouts): layout_idx = len(prs.slide_layouts) - 1
                slide_layout = prs.slide_layouts[layout_idx]
                slide = prs.slides.add_slide(slide_layout)
                return len(prs.slides) - 1
            except Exception as e:
                raise FMS_Error(f"Failed to add slide: {str(e)}", error_type="Runtime Error")
        def pptx_add_content_slide(i, a):
            prs = _get_doc(a[0])
            title = str(a[1])
            bullets = a[2] if len(a) > 2 else []
            try:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                if len(slide.placeholders) > 1 and bullets:
                    tf = slide.placeholders[1].text_frame
                    tf.text = ""
                    for idx, bullet in enumerate(bullets):
                        if idx == 0: tf.paragraphs[0].text = str(bullet)
                        else:
                            p = tf.add_paragraph()
                            p.text = str(bullet)
                return len(prs.slides) - 1
            except Exception as e:
                raise FMS_Error(f"Failed to add content slide: {str(e)}", error_type="Runtime Error")
        def pptx_add_text_box(i, a):
            prs = _get_doc(a[0])
            slide_idx = int(a[1])
            text = str(a[2])
            left = float(a[3]) if len(a) > 3 else 1.0
            top = float(a[4]) if len(a) > 4 else 1.0
            width = float(a[5]) if len(a) > 5 else 4.0
            height = float(a[6]) if len(a) > 6 else 2.0
            try:
                from pptx.util import Inches
                slide = prs.slides[slide_idx]
                txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
                tf = txBox.text_frame
                tf.text = text
                tf.word_wrap = True
                return "SUCCESS"
            except IndexError:
                raise FMS_Error(f"Slide index {slide_idx} out of range", error_type="Runtime Error")
            except Exception as e:
                raise FMS_Error(f"Failed to add text box: {str(e)}", error_type="Runtime Error")
        def pptx_add_image(i, a):
            prs = _get_doc(a[0])
            slide_idx = int(a[1])
            img_path = str(a[2])
            left = float(a[3]) if len(a) > 3 else 1.0
            top = float(a[4]) if len(a) > 4 else 2.0
            width = float(a[5]) if len(a) > 5 else 4.0
            height = float(a[6]) if len(a) > 6 else 3.0
            try:
                from pptx.util import Inches
                slide = prs.slides[slide_idx]
                slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
                return "SUCCESS"
            except FileNotFoundError:
                raise FMS_Error(f"Image not found: '{img_path}'", error_type="File Error")
            except Exception as e:
                raise FMS_Error(f"Failed to add image: {str(e)}", error_type="Runtime Error")
        def pptx_add_table(i, a):
            prs = _get_doc(a[0])
            slide_idx = int(a[1])
            data = a[2]
            left = float(a[3]) if len(a) > 3 else 0.5
            top = float(a[4]) if len(a) > 4 else 2.0
            width = float(a[5]) if len(a) > 5 else 9.0
            height = float(a[6]) if len(a) > 6 else 3.0
            if not isinstance(data, list) or len(data) == 0:
                raise FMS_Error("PPTX_ADD_TABLE requires a non-empty list of lists", error_type="Runtime Error")
            try:
                from pptx.util import Inches
                slide = prs.slides[slide_idx]
                rows = len(data)
                cols = max(len(row) for row in data)
                table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
                table = table_shape.table
                for r, row_data in enumerate(data):
                    for c, cell_val in enumerate(row_data):
                        if c < cols:
                            table.cell(r, c).text = str(cell_val) if cell_val is not None else ""
                return "SUCCESS"
            except Exception as e:
                raise FMS_Error(f"Failed to add table: {str(e)}", error_type="Runtime Error")
        def pptx_get_slide_count(i, a):
            prs = _get_doc(a[0])
            return len(prs.slides)
        def pptx_get_text(i, a):
            prs = _get_doc(a[0])
            slides_text = []
            for idx, slide in enumerate(prs.slides):
                slide_texts = [f"=== Slide {idx + 1} ==="]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip(): slide_texts.append(para.text)
                slides_text.append("\n".join(slide_texts))
            return "\n".join(slides_text)
        def pptx_get_slides_text(i, a):
            prs = _get_doc(a[0])
            result = []
            for slide in prs.slides:
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip(): texts.append(para.text)
                result.append("\n".join(texts))
            return result
        def pptx_delete_slide(i, a):
            prs = _get_doc(a[0])
            slide_idx = int(a[1])
            try:
                rId = prs.slides._sldIdLst[slide_idx].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[slide_idx]
                return "SUCCESS"
            except Exception as e:
                raise FMS_Error(f"Failed to delete slide: {str(e)}", error_type="Runtime Error")
        def pptx_set_slide_bg(i, a):
            prs = _get_doc(a[0])
            slide_idx = int(a[1])
            color_hex = str(a[2]).lstrip('#')
            try:
                pptx = self._get_pptx()
                slide = prs.slides[slide_idx]
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = pptx.dml.color.RGBColor.from_string(color_hex)
                return "SUCCESS"
            except Exception as e:
                raise FMS_Error(f"Failed to set background: {str(e)}", error_type="Runtime Error")
        def pptx_to_pdf(i, a):
            pptx_path, pdf_path = str(a[0]), str(a[1])
            try:
                result = subprocess.run(
                    ['soffice', '--headless', '--convert-to', 'pdf', '--outdir',
                     os.path.dirname(pdf_path) or '.', pptx_path],
                    capture_output=True, text=True, timeout=60
                )
                if result.returncode == 0: return "SUCCESS"
                else: return f"ERROR: LibreOffice conversion failed: {result.stderr}"
            except FileNotFoundError:
                return "ERROR: LibreOffice not found. Install LibreOffice for PPTX to PDF conversion."
            except Exception as e:
                return f"ERROR: {str(e)}"
        def pptx_add_notes(i, a):
            prs = _get_doc(a[0])
            slide_idx = int(a[1])
            notes_text = str(a[2])
            try:
                slide = prs.slides[slide_idx]
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes_text
                return "SUCCESS"
            except Exception as e:
                raise FMS_Error(f"Failed to add notes: {str(e)}", error_type="Runtime Error")

        env.define("PPTX_CREATE", NativeFunction("PPTX_CREATE", 0, pptx_create))
        env.define("PPTX_OPEN", NativeFunction("PPTX_OPEN", 1, pptx_open))
        env.define("PPTX_SAVE", NativeFunction("PPTX_SAVE", 2, pptx_save))
        env.define("PPTX_CLOSE", NativeFunction("PPTX_CLOSE", 1, pptx_close))
        env.define("PPTX_ADD_TITLE_SLIDE", NativeFunction("PPTX_ADD_TITLE_SLIDE", -1, pptx_add_title_slide))
        env.define("PPTX_ADD_SLIDE", NativeFunction("PPTX_ADD_SLIDE", -1, pptx_add_slide))
        env.define("PPTX_ADD_CONTENT_SLIDE", NativeFunction("PPTX_ADD_CONTENT_SLIDE", -1, pptx_add_content_slide))
        env.define("PPTX_ADD_TEXT_BOX", NativeFunction("PPTX_ADD_TEXT_BOX", -1, pptx_add_text_box))
        env.define("PPTX_ADD_IMAGE", NativeFunction("PPTX_ADD_IMAGE", -1, pptx_add_image))
        env.define("PPTX_ADD_TABLE", NativeFunction("PPTX_ADD_TABLE", -1, pptx_add_table))
        env.define("PPTX_GET_SLIDE_COUNT", NativeFunction("PPTX_GET_SLIDE_COUNT", 1, pptx_get_slide_count))
        env.define("PPTX_GET_TEXT", NativeFunction("PPTX_GET_TEXT", 1, pptx_get_text))
        env.define("PPTX_GET_SLIDES_TEXT", NativeFunction("PPTX_GET_SLIDES_TEXT", 1, pptx_get_slides_text))
        env.define("PPTX_DELETE_SLIDE", NativeFunction("PPTX_DELETE_SLIDE", 2, pptx_delete_slide))
        env.define("PPTX_SET_SLIDE_BG", NativeFunction("PPTX_SET_SLIDE_BG", 3, pptx_set_slide_bg))
        env.define("PPTX_TO_PDF", NativeFunction("PPTX_TO_PDF", 2, pptx_to_pdf))
        env.define("PPTX_ADD_NOTES", NativeFunction("PPTX_ADD_NOTES", 3, pptx_add_notes))